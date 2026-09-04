from io import BytesIO
from pathlib import Path
from csv import reader

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

import models


ROOT = Path(__file__).parents[1]
HTML = (ROOT / "templates/index.html").read_text(encoding="utf-8")
JS = (ROOT / "static/js/app.js").read_text(encoding="utf-8")


def _presentation(response):
    return Presentation(BytesIO(response.data))


def _text_shapes(slide):
    return [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]


def _shape_with_text(slide, text):
    return next(shape for shape in _text_shapes(slide) if shape.text == text)


def _shape_font_names(shape):
    names = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.name:
                names.append(run.font.name)
    return names


def _title_color(slide):
    title = slide.shapes[0]
    return title.text_frame.paragraphs[0].runs[0].font.color.rgb


def _cell_color(cell):
    return cell.text_frame.paragraphs[0].runs[0].font.color.rgb


def test_pptx_export_starts_with_overview_table_and_status_colors(client):
    models.create_sample({"id": "ppt-success", "target_product": "Target-S", "status": 1})
    models.create_sample({"id": "ppt-fail", "target_product": "Target-F", "status": 0})
    models.create_sample({"id": "ppt-pending", "target_product": "Target-P", "status": 2})

    presentation = _presentation(client.get("/api/samples/export_pptx?lang=zh"))

    assert len(presentation.slides) == 4
    cover = presentation.slides[0]
    assert cover.shapes[0].text == "样品总览"
    table_shape = next(shape for shape in cover.shapes if shape.has_table)
    table = table_shape.table
    assert len(table.columns) == 3
    assert [cell.text for cell in table.rows[0].cells] == ["编号", "目标样品", "是否成功"]

    rows = {
        table.rows[index].cells[0].text: table.rows[index]
        for index in range(1, len(table.rows))
    }
    assert rows["ppt-success"].cells[1].text == "Target-S"
    assert rows["ppt-success"].cells[2].text == "成功"
    assert rows["ppt-fail"].cells[2].text == "失败"
    assert rows["ppt-pending"].cells[2].text == "待定"
    assert _cell_color(rows["ppt-success"].cells[2]) == RGBColor(22, 163, 74)
    assert _cell_color(rows["ppt-fail"].cells[2]) == RGBColor(220, 38, 38)
    assert _cell_color(rows["ppt-pending"].cells[2]) == RGBColor(0, 0, 0)


def test_pptx_zh_export_has_one_slide_per_sample_and_requested_vertical_regions(client):
    models.create_sample(
        {
            "id": "ppt-zh",
            "target_product": "BiNiTe",
            "status": 1,
            "element_ratios": [
                {"element": "Te", "ratio": 1},
                {"element": "Bi", "ratio": 2},
                {"element": "Ni", "ratio": 1},
            ],
            "growth_process": "中文生长方法",
            "growth_process_en": "English growth method",
            "results": "中文结果文字",
            "results_en": "English result text",
        }
    )

    response = client.get("/api/samples/export_pptx?lang=zh")

    assert response.status_code == 200
    assert response.mimetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    presentation = _presentation(response)
    assert len(presentation.slides) == 2
    slide = presentation.slides[1]
    assert slide.shapes[0].text == "ppt-zh-BiNiTe-成功"
    assert _shape_with_text(slide, "结果").top < _shape_with_text(slide, "生长方法").top
    assert _shape_with_text(slide, "中文结果文字")
    assert _shape_with_text(slide, "元素比例：Bi2NiTe\n中文生长方法")
    growth_box = next(
        shape for shape in slide.shapes
        if not shape.text and shape.top == Inches(4.4)
    )
    assert growth_box.height == Inches(1.25)
    assert "English result text" not in "\n".join(shape.text for shape in _text_shapes(slide))
    assert "English growth method" not in "\n".join(shape.text for shape in _text_shapes(slide))
    assert all(name == "Microsoft YaHei" for shape in _text_shapes(slide) for name in _shape_font_names(shape))
    assert _title_color(slide) == RGBColor(22, 163, 74)


def test_pptx_en_export_uses_english_fields_times_new_roman_and_failure_color(client):
    models.create_sample(
        {
            "id": "ppt-en",
            "target_product": "BiNiTe",
            "status": 0,
            "element_ratios": [{"element": "Bi", "ratio": 1}],
            "growth_process": "中文生长方法",
            "growth_process_en": "English growth method",
            "results": "中文结果文字",
            "results_en": "English result text",
        }
    )

    response = client.get("/api/samples/export_pptx?lang=en")

    assert response.status_code == 200
    presentation = _presentation(response)
    slide = presentation.slides[1]
    assert slide.shapes[0].text == "ppt-en-BiNiTe-Fail"
    assert _shape_with_text(slide, "Results")
    assert _shape_with_text(slide, "Growth Method")
    assert _shape_with_text(slide, "English result text")
    assert _shape_with_text(slide, "Element Ratio: Bi\nEnglish growth method")
    all_text = "\n".join(shape.text for shape in _text_shapes(slide))
    assert "中文结果文字" not in all_text
    assert "中文生长方法" not in all_text
    assert all(name == "Times New Roman" for shape in _text_shapes(slide) for name in _shape_font_names(shape))
    assert _title_color(slide) == RGBColor(220, 38, 38)


@pytest.mark.parametrize(
    ("status", "expected_zh", "expected_en", "color"),
    [
        (2, "待定", "Pending", RGBColor(0, 0, 0)),
        (3, "生长中", "Growing", RGBColor(0, 0, 0)),
        (4, "生长完成", "Done", RGBColor(0, 0, 0)),
    ],
)
def test_pptx_other_status_titles_are_black(client, status, expected_zh, expected_en, color):
    sample_id = f"ppt-status-{status}"
    models.create_sample({"id": sample_id, "target_product": "Target", "status": status})

    zh_slide = _presentation(client.get("/api/samples/export_pptx?lang=zh")).slides[1]
    en_slide = _presentation(client.get("/api/samples/export_pptx?lang=en")).slides[1]

    assert zh_slide.shapes[0].text == f"{sample_id}-Target-{expected_zh}"
    assert en_slide.shapes[0].text == f"{sample_id}-Target-{expected_en}"
    assert _title_color(zh_slide) == color
    assert _title_color(en_slide) == color


def test_pptx_export_keeps_missing_target_language_text_empty_without_fallback(client):
    models.create_sample(
        {
            "id": "ppt-empty-en",
            "target_product": "Target",
            "status": 2,
            "growth_process": "仅中文生长方法",
            "growth_process_en": "",
            "results": "仅中文结果",
            "results_en": "",
        }
    )

    slide = _presentation(client.get("/api/samples/export_pptx?lang=en")).slides[1]
    all_text = "\n".join(shape.text for shape in _text_shapes(slide))
    assert "仅中文生长方法" not in all_text
    assert "仅中文结果" not in all_text
    assert not any(shape.text == "Element Ratio: " for shape in _text_shapes(slide))


def test_pptx_export_button_and_route_are_exposed_in_both_languages():
    assert 'id="exportPptxBtn"' in HTML
    assert 'onclick="exportSamplesPptx()"' in HTML
    assert "exportPptx" in JS
    assert "/api/samples/export_pptx?lang=${currentLang}" in JS
    assert "format === 'pptx'" in JS
    assert "导出 PPTX" in JS
    assert "Export PPTX" in JS
