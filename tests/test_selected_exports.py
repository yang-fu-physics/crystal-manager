from csv import reader
from hashlib import sha256
from io import BytesIO, StringIO

import pytest
from pptx import Presentation

import models


def _csv_rows(response):
    return list(reader(StringIO(response.data.decode("utf-8-sig"))))


def _sample(sample_id, **overrides):
    data = {
        "id": sample_id,
        "target_product": f"Target-{sample_id}",
        "status": 1,
        "growth_process": f"Chinese process {sample_id}",
        "growth_process_en": f"English process {sample_id}",
        "results": f"Chinese result {sample_id}",
        "results_en": f"English result {sample_id}",
    }
    data.update(overrides)
    return data


def _db_snapshot():
    conn = models.get_db()
    try:
        rows = conn.execute(
            "SELECT id, target_product, is_successful, growth_process, "
            "growth_process_en, results, results_en FROM samples ORDER BY id"
        ).fetchall()
    finally:
        conn.close()
    return sha256(repr([tuple(row) for row in rows]).encode("utf-8")).hexdigest()


def test_post_csv_exports_only_selected_samples_in_reverse_current_order(client):
    models.create_sample(_sample("sample-A", sintering_start="2024-01-01T00:00"))
    models.create_sample(_sample("sample-B", sintering_start="2024-01-02T00:00"))
    models.create_sample(_sample("sample-C", sintering_start="2024-01-03T00:00"))
    current_order = [sample["id"] for sample in models.get_all_samples()]
    selected_ids = current_order[:2]

    response = client.post(
        "/api/samples/export?lang=zh",
        json={"sample_ids": selected_ids},
    )

    assert response.status_code == 200
    rows = _csv_rows(response)
    assert [row[0] for row in rows[1:]] == list(reversed(selected_ids))
    assert len(rows) == 3


def test_post_pptx_exports_only_selected_samples_in_reverse_current_order(client):
    models.create_sample(_sample("sample-A", sintering_start="2024-01-01T00:00"))
    models.create_sample(_sample("sample-B", sintering_start="2024-01-02T00:00"))
    models.create_sample(_sample("sample-C", sintering_start="2024-01-03T00:00"))
    current_order = models.get_all_samples()
    selected_samples = current_order[:2]
    selected_ids = [sample["id"] for sample in selected_samples]

    response = client.post(
        "/api/samples/export_pptx?lang=en",
        json={"sample_ids": selected_ids},
    )

    assert response.status_code == 200
    presentation = Presentation(BytesIO(response.data))
    assert len(presentation.slides) == 3
    expected_titles = [
        f"{sample['id']}-{sample.get('target_product') or ''}-Success"
        for sample in reversed(selected_samples)
    ]
    assert [slide.shapes[0].text for slide in list(presentation.slides)[1:]] == expected_titles


def test_selected_exports_keep_language_specific_fields(client):
    models.create_sample(
        _sample(
            "sample-language",
            results="Chinese only",
            results_en="English only",
            growth_process="Chinese growth only",
            growth_process_en="English growth only",
        )
    )

    zh_row = _csv_rows(
        client.post(
            "/api/samples/export?lang=zh",
            json={"sample_ids": ["sample-language"]},
        )
    )[1]
    en_row = _csv_rows(
        client.post(
            "/api/samples/export?lang=en",
            json={"sample_ids": ["sample-language"]},
        )
    )[1]

    assert zh_row[3] == "Chinese only"
    assert "Chinese growth only" in zh_row[4]
    assert en_row[3] == "English only"
    assert "English growth only" in en_row[4]
    assert "Chinese only" not in en_row
    assert "Chinese growth only" not in en_row


def test_get_exports_reverse_current_order_for_all_samples(client):
    models.create_sample(_sample("sample-A", sintering_start="2024-01-01T00:00"))
    models.create_sample(_sample("sample-B", sintering_start="2024-01-02T00:00"))
    models.create_sample(_sample("sample-C", sintering_start="2024-01-03T00:00"))
    current_order = models.get_all_samples()
    expected_ids = [sample["id"] for sample in reversed(current_order)]

    csv_response = client.get("/api/samples/export?lang=zh")
    pptx_response = client.get("/api/samples/export_pptx?lang=en")

    assert csv_response.status_code == 200
    assert len(_csv_rows(csv_response)) == 4
    assert pptx_response.status_code == 200
    assert len(Presentation(BytesIO(pptx_response.data)).slides) == 4
    assert [row[0] for row in _csv_rows(csv_response)[1:]] == expected_ids
    sample_by_id = {sample["id"]: sample for sample in current_order}
    expected_titles = [
        f"{sample_by_id[sample_id]['id']}-{sample_by_id[sample_id].get('target_product') or ''}-Success"
        for sample_id in expected_ids
    ]
    presentation = Presentation(BytesIO(pptx_response.data))
    assert [slide.shapes[0].text for slide in list(presentation.slides)[1:]] == expected_titles


def test_duplicate_selected_ids_are_deduplicated_for_reverse_current_order(client):
    models.create_sample(_sample("sample-A"))
    models.create_sample(_sample("sample-B"))

    response = client.post(
        "/api/samples/export?lang=zh",
        json={"sample_ids": ["sample-B", "sample-B", "sample-A"]},
    )

    assert response.status_code == 200
    assert [row[0] for row in _csv_rows(response)[1:]] == ["sample-B", "sample-A"]


@pytest.mark.parametrize("path", ["/api/samples/export", "/api/samples/export_pptx"])
@pytest.mark.parametrize(
    "payload",
    [
        None,
        {},
        {"sample_ids": "sample-A"},
        {"sample_ids": []},
        {"sample_ids": ["sample-A", 1]},
        {"sample_ids": ["missing-sample"]},
    ],
)
def test_invalid_selected_export_requests_return_400(client, path, payload):
    models.create_sample(_sample("sample-A"))

    if payload is None:
        response = client.post(path + "?lang=zh")
    else:
        response = client.post(path + "?lang=zh", json=payload)

    assert response.status_code == 400


def test_selected_exports_do_not_write_to_database(client):
    models.create_sample(_sample("sample-A"))
    models.create_sample(_sample("sample-B"))
    before = _db_snapshot()

    csv_response = client.post(
        "/api/samples/export?lang=zh",
        json={"sample_ids": ["sample-B"]},
    )
    pptx_response = client.post(
        "/api/samples/export_pptx?lang=en",
        json={"sample_ids": ["sample-A"]},
    )

    assert csv_response.status_code == 200
    assert pptx_response.status_code == 200
    assert _db_snapshot() == before
